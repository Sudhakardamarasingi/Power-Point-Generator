import io
import json
import requests
import streamlit as st
from pptx import Presentation

# ---------- CONFIG ----------
WEBHOOK_URL = "https://sudha-mad-max-1997.app.n8n.cloud/webhook/f4892281-e1a0-429c-ae0a-16661a18e576"

st.set_page_config(
    page_title="Agentic PPT Designer",
    page_icon="🎨",
    layout="wide",
)


# ---------- STYLES (Canva-ish) ----------
st.markdown(
    """
    <style>
    /* App background */
    .stApp {
        background: radial-gradient(circle at top left, #2b3a67 0, #050816 45%, #050816 100%);
        color: #f9fafb;
        font-family: system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }

    /* Center main content */
    .main-block {
        max-width: 980px;
        margin: 0 auto;
        padding-top: 2.2rem;
        padding-bottom: 3rem;
    }

    /* Card container */
    .glass-card {
        background: rgba(15, 23, 42, 0.88);
        border-radius: 18px;
        padding: 1.75rem 2rem;
        border: 1px solid rgba(148, 163, 184, 0.35);
        box-shadow: 0 18px 60px rgba(15, 23, 42, 0.75);
        backdrop-filter: blur(16px);
    }

    .pill-badge {
        display: inline-flex;
        align-items: center;
        gap: 0.35rem;
        border-radius: 999px;
        border: 1px solid rgba(148, 163, 184, 0.55);
        padding: 0.2rem 0.65rem;
        font-size: 0.75rem;
        opacity: 0.95;
    }

    .pill-dot {
        width: 8px;
        height: 8px;
        border-radius: 999px;
        background: linear-gradient(135deg,#22c55e,#a855f7);
    }

    .topic-chip {
        display: inline-block;
        padding: 0.25rem 0.65rem;
        border-radius: 999px;
        font-size: 0.8rem;
        margin: 0.15rem 0.3rem 0.15rem 0;
        background: rgba(15, 23, 42, 0.8);
        border: 1px solid rgba(148, 163, 184, 0.5);
        cursor: default;
    }

    .progress-label {
        font-size: 0.85rem;
        opacity: 0.9;
        margin-top: 0.4rem;
    }

    /* Make textarea look nicer */
    textarea {
        border-radius: 14px !important;
        border: 1px solid rgba(148, 163, 184, 0.6) !important;
    }

    /* Buttons */
    button[kind="primary"] {
        border-radius: 999px !important;
        font-weight: 600 !important;
    }
    </style>
    """,
    unsafe_allow_html=True,
)


# ---------- HELPERS ----------
def parse_ai_response(response: requests.Response) -> dict:
    """
    n8n Respond to Webhook returns:
    - plain text body that is already a JSON string, from {{$json["output"]}}
    """
    raw = response.text.strip()

    if not raw:
        raise ValueError("AI response is empty")

    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        # Show what we got so we can debug if needed
        st.error(f"Could not parse AI JSON response: {e}")
        st.code(raw, language="json")
        raise


def build_ppt_from_spec(spec: dict) -> bytes:
    """
    Build a .pptx from the JSON spec.

    Expected format:
    {
      "title": "Presentation Title",
      "slides": [
        {
          "heading": "Slide heading",
          "bullets": ["point 1", "point 2"],
          "notes": "optional presenter notes",
          "image_prompt": "optional"
        },
        ...
      ]
    }
    """
    prs = Presentation()

    slides = spec.get("slides", [])
    if not slides:
        raise ValueError("No 'slides' array found in AI response")

    # Title slide
    first = slides[0]
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = spec.get("title") or first.get("heading", "Presentation")
    subtitle_shape.text = first.get("notes", "")

    # Content slides
    content_layout = prs.slide_layouts[1]

    for slide_spec in slides[1:]:
        slide = prs.slides.add_slide(content_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_shape.text = slide_spec.get("heading", "")

        text_frame = body_shape.text_frame
        text_frame.clear()

        bullets = slide_spec.get("bullets", [])
        if bullets:
            text_frame.text = bullets[0]
            for bullet in bullets[1:]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ---------- STATE ----------
if "ppt_bytes" not in st.session_state:
    st.session_state["ppt_bytes"] = None


# ---------- HEADER ----------
st.markdown(
    """
    <div class="main-block">
      <div style="text-align:center; margin-bottom:1.2rem;">
        <div class="pill-badge">
          <span class="pill-dot"></span>
          <span>Agentic AI + n8n + Streamlit</span>
        </div>
        <h1 style="margin-top:1rem; margin-bottom:0.4rem;">
          Agentic PowerPoint Designer
        </h1>
        <p style="opacity:0.85; max-width:620px; margin:0 auto;">
          Turn a single prompt into a complete, editable slide deck.
          Powered by a Gemini agent orchestrated in n8n.
        </p>
      </div>
    """,
    unsafe_allow_html=True,
)


# ---------- MAIN CARD ----------
with st.container():
    col_left, col_right = st.columns([1.25, 1])

    with col_left:
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)

        st.subheader("Describe your presentation", anchor=False)
        st.caption("Be clear about topic, audience, and depth. The agent will design the slides for you.")

        prompt = st.text_area(
            label="",
            placeholder="Example: Create an 8-slide PPT explaining agentic AI for non-technical business leaders...",
            height=160,
        )

        st.markdown(
            """
            <div style="margin-top:0.3rem; margin-bottom:0.6rem;">
              <span class="topic-chip">Intro to Agentic AI</span>
              <span class="topic-chip">ML model lifecycle for churn prediction</span>
              <span class="topic-chip">GenAI use cases for marketing</span>
              <span class="topic-chip">Cloud vs On-prem analytics</span>
            </div>
            """,
            unsafe_allow_html=True,
        )

        generate = st.button("✨ Generate PPT", use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

    with col_right:
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        st.subheader("What this tool does", anchor=False)
        st.markdown(
            """
            - Uses an **agentic workflow** in n8n to call Gemini.  
            - Gemini outputs a structured **JSON slide outline**.  
            - This app converts that JSON into a real `.pptx` file.  
            - You can edit the deck in PowerPoint, Google Slides, or Keynote.
            """
        )
        st.markdown("---")
        st.markdown(
            """
            **Tips for best results**
            - Specify audience: *students, managers, data team, CXOs*  
            - Mention level: *beginner / intermediate / advanced*  
            - Add constraints: *“no more than 5 bullets per slide”*,  
              *“include examples and use cases”*, etc.
            """
        )
        st.markdown("</div>", unsafe_allow_html=True)

st.markdown("</div>", unsafe_allow_html=True)  # close main-block


# ---------- ACTION + PROGRESS ----------
if generate:
    if not prompt or not prompt.strip():
        st.warning("Please enter a topic or description for the PPT.")
    else:
        progress = st.progress(0)
        status_text = st.empty()

        try:
            status_text.markdown("🔗 Connecting to agentic backend...")
            progress.progress(15)

            resp = requests.post(
                WEBHOOK_URL,
                json={"prompt": prompt},
                timeout=180,
            )

            progress.progress(45)
            status_text.markdown("🤖 Agent is designing the slide structure...")

            if resp.status_code != 200:
                progress.progress(0)
                st.error(
                    f"PPT generation failed.\n"
                    f"Status: {resp.status_code}\n"
                    f"Body: {resp.text}"
                )
            else:
                spec = parse_ai_response(resp)

                progress.progress(75)
                status_text.markdown("📑 Converting JSON into a PowerPoint file...")

                ppt_bytes = build_ppt_from_spec(spec)
                st.session_state["ppt_bytes"] = ppt_bytes

                progress.progress(100)
                status_text.markdown("✅ Deck ready! Scroll down to download.")
                size_kb = len(ppt_bytes) / 1024
                st.success(f"PPT generated successfully (~{size_kb:.1f} KB).")

        except Exception as e:
            progress.progress(0)
            st.error(f"Error while generating PPT:\n{e}")


# ---------- DOWNLOAD SECTION ----------
if st.session_state["ppt_bytes"]:
    st.markdown("---")
    st.subheader("Download your presentation", anchor=False)
    st.caption("You can edit this file further in PowerPoint, Google Slides, or Keynote.")

    st.download_button(
        label="⬇️ Download PPTX",
        data=st.session_state["ppt_bytes"],
        file_name="presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
    )
